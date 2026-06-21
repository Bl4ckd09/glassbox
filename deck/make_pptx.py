#!/usr/bin/env python3
"""Generate GlassBox.pptx — a native, editable pitch deck (16:9, dark theme)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette
BG    = RGBColor(0x08, 0x08, 0x0B)
PANEL = RGBColor(0x16, 0x16, 0x1B)
PANEL2= RGBColor(0x24, 0x24, 0x2E)
WHITE = RGBColor(0xF4, 0xF4, 0xF5)
ZINC_L= RGBColor(0xD4, 0xD4, 0xD8)
ZINC  = RGBColor(0xA1, 0xA1, 0xAA)
MUTE  = RGBColor(0x71, 0x71, 0x7A)
EMER  = RGBColor(0x34, 0xD3, 0x99)
SKY   = RGBColor(0x38, 0xBD, 0xF8)
VIOLET= RGBColor(0xA7, 0x8B, 0xFA)
ROSE  = RGBColor(0xFB, 0x71, 0x85)

BODY = "Arial"
MONO = "Menlo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font) tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        if isinstance(para, tuple): para = [para]
        for (txt, size, color, bold, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb

def chip(s, x, y, label, color, w=1.4, h=0.42):
    box(s, x, y, w, h, fill=PANEL2)
    text(s, x, y, w, h, [(label, 13, color, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x + w + 0.15

def kicker(s, label, color):
    text(s, 0.9, 0.6, 11, 0.4, [(label, 15, color, True, BODY)])

def footer(s):
    text(s, 0.9, 7.02, 11.5, 0.35,
         [[("glassbox-beige.vercel.app", 11, MUTE, False, BODY),
           ("    ·    github.com/Bl4ckd09/glassbox", 11, MUTE, False, BODY)]])

# ---------------- 1 · Title ----------------
s = slide()
box(s, 0.9, 2.35, 0.55, 0.55, fill=EMER)
text(s, 0.9, 2.35, 0.55, 0.55, [("G", 26, BG, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.6, 2.32, 6, 0.6, [("GlassBox", 30, WHITE, True, BODY)], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 3.15, 11.5, 1.4, [("Pre-registration for trading calls.", 46, WHITE, True, BODY)])
text(s, 0.9, 4.25, 11.5, 0.7, [("Audit, don't trust.", 26, EMER, True, BODY)])
x = 0.9
for label, c in [("Sui", SKY), ("DeepBook", EMER), ("Walrus", VIOLET), ("AI-native", ZINC_L)]:
    x = chip(s, x, 5.15, label, c, w=1.25)
text(s, 0.9, 6.2, 11.5, 0.5,
     [("Tamper-proof, auditable track records on Sui — for the Encode Vibe Coding Hackathon.", 15, ZINC, False, BODY)])
footer(s)

# ---------------- 2 · Problem ----------------
s = slide(); kicker(s, "THE PROBLEM", ROSE)
text(s, 0.9, 1.15, 11.5, 1.0, [("Finfluencers delete their losing calls.", 40, WHITE, True, BODY)])
text(s, 0.9, 2.5, 11.0, 2.5, [
    [("They post 50 trades, quietly remove the 30 that lost, screenshot the 20 that won, and sell you a ", 22, ZINC_L, False, BODY),
     ("100%-win-rate fantasy.", 22, WHITE, True, BODY)],
    ("", 8, ZINC, False, BODY),
    ("Retail traders can't tell real skill from survivorship bias — a problem the UK FCA and US SEC have both flagged about misleading trading promotions.", 22, ZINC, False, BODY),
])
box(s, 0.9, 5.3, 11.53, 1.05, fill=PANEL)
text(s, 1.2, 5.3, 11.0, 1.05,
     [[("This is an ", 18, ZINC_L, False, BODY), ("information-asymmetry", 18, EMER, True, BODY),
       (" problem — retail is structurally disadvantaged versus whoever controls the record.", 18, ZINC_L, False, BODY)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ---------------- 3 · Insight ----------------
s = slide(); kicker(s, "THE INSIGHT", SKY)
text(s, 0.9, 1.3, 11.5, 0.8, [("On-chain PnL trackers show what a wallet ", 28, ZINC, False, BODY)])
text(s, 0.9, 1.85, 11.5, 0.8, [[("did", 28, ZINC, True, BODY), (".", 28, ZINC, False, BODY)]])
text(s, 0.9, 2.9, 11.7, 1.6, [
    [("GlassBox proves what a strategist ", 34, WHITE, True, BODY), ("said they'd do", 34, WHITE, True, BODY)],
    [("— in advance, with the receipts.", 34, EMER, True, BODY)],
])
box(s, 0.9, 5.2, 11.53, 1.1, fill=PANEL)
text(s, 1.2, 5.2, 11.0, 1.1,
     [("It's clinical-trial pre-registration, applied to trading.", 24, ZINC_L, True, BODY)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ---------------- 4 · How it works ----------------
s = slide(); kicker(s, "HOW IT WORKS", EMER)
steps = [
    ("1", "Computed on real markets", "Built from live Sui DeepBook order books and trade history — a real on-chain CLOB, not a simulation.", SKY),
    ("2", "Pre-registered before the outcome", "The full call — side, reasoning, risk plan — is signed with a Sui key and written to Walrus the instant it's made. Content-addressed, so it can't be edited, deleted, or back-dated.", EMER),
    ("3", "Verifiable by anyone", "Re-fetch the blob from a public Walrus aggregator and check the signature. No need to trust GlassBox.", VIOLET),
]
y = 1.35
for n, t, body, c in steps:
    box(s, 0.9, y, 11.53, 1.62, fill=PANEL)
    box(s, 1.15, y+0.5, 0.62, 0.62, fill=c)
    text(s, 1.15, y+0.5, 0.62, 0.62, [(n, 22, BG, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.0, y+0.22, 10.2, 0.5, [(t, 22, WHITE, True, BODY)])
    text(s, 2.0, y+0.72, 10.2, 0.8, [(body, 15, ZINC, False, BODY)])
    y += 1.78
footer(s)

# ---------------- 5 · Live proof ----------------
s = slide(); kicker(s, "A REAL CALL — VERIFIED LIVE", EMER)
box(s, 0.9, 1.25, 11.53, 1.85, fill=PANEL)
text(s, 1.2, 1.45, 3, 0.7, [("HOLD", 30, ZINC, True, BODY)])
text(s, 2.9, 1.5, 8, 0.6, [("SUI / USDC", 24, WHITE, True, BODY)])
text(s, 2.9, 2.05, 9, 0.5, [("entry 0.71127   ·   conviction -0.09   ·   transparent weighted-vote engine", 14, MUTE, False, BODY)])
x = 1.2
for lbl, c in [("RSI 63.2", ZINC_L), ("stop 1.5%", ROSE), ("target 2.7%", EMER), ("size 0.88%", ZINC_L), ("signed", VIOLET), ("anchored on Walrus", EMER)]:
    x = chip(s, x, 2.55, lbl, c, w=1.55 if len(lbl) > 9 else 1.15)
box(s, 0.9, 3.35, 11.53, 3.05, fill=RGBColor(0x06,0x06,0x08), line=PANEL2, line_w=1.5)
text(s, 1.2, 3.55, 10, 0.4, [("IMMUTABLE PROOF (Walrus testnet)", 14, MUTE, True, BODY)])
text(s, 1.2, 4.05, 11, 0.5, [("blobId:  77VWcMAKjI…I_6qpyfo", 18, ZINC_L, False, MONO)])
text(s, 1.2, 4.55, 11, 0.4, [("pre-registered before outcome · Walrus epoch 435 · Sui obj 0xd300…130b", 13, MUTE, False, BODY)])
text(s, 1.2, 5.15, 11, 0.5, [[("✓  ", 20, EMER, True, BODY), ("Fetched from public Walrus aggregator — content matches.", 20, EMER, True, BODY)]])
text(s, 1.2, 5.7, 11, 0.5, [[("✓  ", 20, EMER, True, BODY), ("Signature valid — authored by 0xcb37…f673.", 20, EMER, True, BODY)]])
footer(s)

# ---------------- 6 · Leaderboard ----------------
s = slide(); kicker(s, "VERIFIABLE LEADERBOARD", SKY)
text(s, 0.9, 1.05, 11.5, 0.5, [("Ranked by a disclosed, risk-adjusted GlassBox Score — not raw PnL.", 17, ZINC, False, BODY)])
rows = [("1", "Diamond Hands Dan", "", "78.7", "+3.80%", "83%", EMER),
        ("2", "Momentum Mia", "", "65.1", "+1.12%", "57%", EMER),
        ("3", "Aletheia", "AI", "49.6", "-0.58%", "33%", ROSE)]
# header
text(s, 8.4, 1.62, 1.6, 0.3, [("SCORE", 11, MUTE, True, BODY)])
text(s, 9.9, 1.62, 1.6, 0.3, [("RETURN", 11, MUTE, True, BODY)])
text(s, 11.3, 1.62, 1.2, 0.3, [("WIN", 11, MUTE, True, BODY)])
y = 1.95
for n, name, tag, score, ret, win, rc in rows:
    box(s, 0.9, y, 11.53, 1.25, fill=PANEL)
    text(s, 1.1, y, 0.5, 1.25, [(n, 22, MUTE, True, BODY)], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 1.65, y+0.33, 0.6, 0.6, fill=SKY)
    ini = "".join(w[0] for w in name.split())[:2]
    text(s, 1.65, y+0.33, 0.6, 0.6, [(ini, 18, BG, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    nm = [(name, 22, WHITE, True, BODY)]
    if tag: nm.append(("   "+tag, 13, SKY, True, BODY))
    text(s, 2.45, y+0.22, 5.6, 0.5, [nm])
    text(s, 2.45, y+0.74, 5.6, 0.35, [[("✓ ", 12, EMER, True, BODY), ("100% on Walrus", 12, EMER, False, BODY)]])
    text(s, 8.4, y, 1.5, 1.25, [(score, 26, SKY, True, BODY)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.9, y, 1.5, 1.25, [(ret, 18, rc, True, BODY)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 11.3, y, 1.2, 1.25, [(win, 18, ZINC_L, True, BODY)], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.4
text(s, 0.9, 6.35, 11.5, 0.4, [("Even our own AI agent (Aletheia) is shown underperforming. Nothing is hidden.", 16, ZINC_L, True, BODY)])
footer(s)

# ---------------- 7 · Glass box, not black box ----------------
s = slide(); kicker(s, "A GLASS BOX, NOT A BLACK BOX", VIOLET)
text(s, 0.9, 1.15, 11.5, 0.8, [("Transparent by construction", 32, WHITE, True, BODY)])
cards = [
    ("Disclosed logic", "Decisions come from a transparent weighted-vote engine — SMA trend, RSI, momentum — every indicator value shown. AI narration explains, never overrides.", SKY),
    ("Built-in risk guards", "Volatility-scaled sizing, hard-capped stops, ~1.8:1 reward:risk, and a spread/liquidity abstain rule that refuses unfair fills.", EMER),
    ("Honest scoring", "The GlassBox Score rewards risk-adjusted quality and penalises recklessness — implementing BGA's 'better systems, not biggest gamble' in the product.", VIOLET),
]
x = 0.9
for t, b, c in cards:
    box(s, x, 2.3, 3.71, 3.7, fill=PANEL)
    text(s, x+0.3, 2.6, 3.1, 0.6, [(t, 19, c, True, BODY)])
    text(s, x+0.3, 3.3, 3.1, 2.5, [(b, 14, ZINC, False, BODY)])
    x += 3.91
footer(s)

# ---------------- 8 · Architecture ----------------
s = slide(); kicker(s, "ARCHITECTURE", SKY)
text(s, 0.9, 1.15, 11.5, 0.7, [("Real on-chain, not mocked", 32, WHITE, True, BODY)])
flow = [
    ("Live Sui DeepBook", "order books + trade history (mainnet CLOB)", SKY),
    ("Transparent engine", "indicators → signal + risk guards + Sui-key signature", EMER),
    ("Walrus", "immutable blob written before the outcome → blobId", VIOLET),
    ("Verify", "anyone re-fetches + checks signature", ZINC_L),
]
x = 0.9
for t, b, c in flow:
    box(s, x, 2.2, 2.75, 2.0, fill=PANEL, line=c, line_w=1.25)
    text(s, x+0.2, 2.45, 2.35, 0.8, [(t, 16, c, True, BODY)])
    text(s, x+0.2, 3.2, 2.35, 0.9, [(b, 12, ZINC, False, BODY)])
    if x < 10:
        text(s, x+2.68, 2.2, 0.35, 2.0, [("→", 22, MUTE, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += 3.04
box(s, 0.9, 4.7, 11.53, 1.5, fill=PANEL)
text(s, 1.2, 4.9, 11, 1.2, [
    [("Stack:  ", 15, WHITE, True, BODY),
     ("Next.js 16 · TypeScript · Tailwind · @mysten/sui · @mysten/deepbook-v3 · @mysten/walrus · Vercel.", 15, ZINC_L, False, BODY)],
    [("Proof of real integration:  ", 15, WHITE, True, BODY),
     ("31 historical calls pre-anchored as real Walrus blobs; live signal generation signs + anchors on demand.", 15, ZINC_L, False, BODY)],
], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ---------------- 9 · Why it wins ----------------
s = slide(); kicker(s, "WHY IT WINS — TWO BOUNTIES, ONE BUILD", EMER)
text(s, 0.9, 1.05, 11.5, 0.6, [("The Walrus audit trail scores on both rubrics at once.", 22, WHITE, True, BODY)])
box(s, 0.9, 1.85, 5.65, 4.4, fill=PANEL)
text(s, 1.2, 2.05, 5.1, 0.5, [("Sui — DeepBook & Walrus", 18, SKY, True, BODY)])
text(s, 1.2, 2.7, 5.1, 3.4, [
    ("Real-World App (50%) — solves a regulator-flagged trust problem.", 14, ZINC_L, False, BODY),
    ("Technical (20%) — real DeepBook reads + real Walrus writes.", 14, ZINC_L, False, BODY),
    ("Product/UX (20%) — clean retail dashboard, live verify.", 14, ZINC_L, False, BODY),
    ("Vision (10%) — a provenance layer for every signal on Sui.", 14, ZINC_L, False, BODY),
])
box(s, 6.78, 1.85, 5.65, 4.4, fill=PANEL)
text(s, 7.08, 2.05, 5.1, 0.5, [("BGA — AI Trading & Strategy", 18, EMER, True, BODY)])
text(s, 7.08, 2.7, 5.1, 3.4, [
    ("Transparency & verifiability (15) — the core feature.", 14, ZINC_L, False, BODY),
    ("Ethos (20) — ranks by risk-adjusted score, not raw PnL.", 14, ZINC_L, False, BODY),
    ("Innovation (20) — pre-registration on decentralized storage.", 14, ZINC_L, False, BODY),
    ("Strategy & risk (15) — disclosed logic + risk guards.", 14, ZINC_L, False, BODY),
    ("Accessibility (10) — plain-language, retail-first.", 14, ZINC_L, False, BODY),
])
footer(s)

# ---------------- 10 · Status / close ----------------
s = slide()
box(s, 5.4, 1.35, 0.75, 0.75, fill=EMER)
text(s, 5.4, 1.35, 0.75, 0.75, [("G", 30, BG, True, BODY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 2.35, 11.5, 0.9, [("GlassBox — shipped & live", 38, WHITE, True, BODY)], align=PP_ALIGN.CENTER)
text(s, 0.9, 3.35, 11.5, 0.5, [("Trading you can audit, not just trust.", 20, EMER, True, BODY)], align=PP_ALIGN.CENTER)
text(s, 0.9, 4.25, 11.5, 0.5, [("glassbox-beige.vercel.app", 22, SKY, True, BODY)], align=PP_ALIGN.CENTER)
text(s, 0.9, 4.85, 11.5, 0.4, [("github.com/Bl4ckd09/glassbox", 16, ZINC, False, BODY)], align=PP_ALIGN.CENTER)
text(s, 0.9, 5.6, 11.5, 0.5, [("Next: on-chain DeepBook execution · Walrus SDK + Sui attestation object · multi-day backtests", 13, MUTE, False, BODY)], align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(__file__), "..", "GlassBox.pptx")
prs.save(out)
print("saved", os.path.abspath(out), "—", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
